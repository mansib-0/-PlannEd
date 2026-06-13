import collections 
try:
    import collections.abc
    collections.Container = collections.abc.Container
    collections.Hashable = collections.abc.Hashable
    collections.Iterable = collections.abc.Iterable
    collections.Iterator = collections.abc.Iterator
    collections.Sized = collections.abc.Sized
    collections.Callable = collections.abc.Callable
    collections.Sequence = collections.abc.Sequence
    collections.Mapping = collections.abc.Mapping
    collections.MutableMapping = collections.abc.MutableMapping
except AttributeError:
    pass

import math
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_poster():
    prs = Presentation()
    
    # A0 Portrait dimensions
    prs.slide_width = Cm(84.1)
    prs.slide_height = Cm(118.9)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Colors
    accent_color = RGBColor(239, 122, 33) # #EF7A21
    text_color = RGBColor(34, 34, 34)
    
    # Header Logo Area
    txBox = slide.shapes.add_textbox(Cm(4), Cm(3), Cm(15), Cm(6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "P·E"
    p.font.bold = True
    p.font.size = Pt(100)
    p.font.name = 'Open Sans'
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Computer Science"
    p2.font.bold = True
    p2.font.size = Pt(26)
    p2.alignment = PP_ALIGN.CENTER
    
    # Bottom border for logo
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4), Cm(9), Cm(15), Cm(0.2)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = text_color
    slide.shapes[-1].line.color.rgb = text_color

    # Title Area
    txBox2 = slide.shapes.add_textbox(Cm(22), Cm(3), Cm(58), Cm(6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "PlannEd: Hands-Free Task Scheduling & Padding"
    p.font.bold = True
    p.font.size = Pt(80)
    p.font.color.rgb = accent_color
    
    p2 = tf2.add_paragraph()
    p2.text = "Researchers: Md. Saiful Islam Mansib & Taniha shiaree tripura"
    p2.font.bold = True
    p2.font.size = Pt(38)
    p2.font.color.rgb = text_color
    p2.space_before = Pt(20)
    
    p3 = tf2.add_paragraph()
    p3.text = "IDs: 2111712042 & 213882042"
    p3.font.bold = True
    p3.font.size = Pt(38)
    p3.font.color.rgb = text_color

    # Orange Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(12), Cm(84.1), Cm(5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = accent_color
    banner.line.color.rgb = accent_color
    
    tf_banner = banner.text_frame
    p_banner = tf_banner.paragraphs[0]
    p_banner.text = "AI-Driven Mobile Application & Machine Learning"
    p_banner.font.bold = True
    p_banner.font.size = Pt(65)
    p_banner.font.color.rgb = RGBColor(255, 255, 255)
    p_banner.alignment = PP_ALIGN.CENTER

    def add_section(x, y, title, text, include_img=None, img_caption=None):
        width = Cm(23)
        # Add Title
        title_box = slide.shapes.add_textbox(x, y, width, Cm(3))
        tf_title = title_box.text_frame
        p_t = tf_title.paragraphs[0]
        p_t.text = title
        p_t.font.bold = True
        p_t.font.size = Pt(52)
        p_t.font.color.rgb = text_color
        
        # Add Title Underline
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Cm(2.8), width, Cm(0.2)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = accent_color
        slide.shapes[-1].line.color.rgb = accent_color
        
        current_y = y + Cm(4)
        
        # Add Text
        if text:
            content_box = slide.shapes.add_textbox(x, current_y, width, Cm(10))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            paragraphs = text.split('\n')
            for i, para_text in enumerate(paragraphs):
                p_c = tf_content.paragraphs[i] if i == 0 else tf_content.add_paragraph()
                p_c.text = para_text
                p_c.font.size = Pt(26)
                p_c.font.color.rgb = text_color
                p_c.space_after = Pt(20)
            
            num_chars = len(text)
            lines = math.ceil(num_chars / 35.0) # approx 35 chars per line at 26pt in 23cm
            text_height = lines * 1.5 
            current_y += Cm(text_height + 1)
            
        if include_img:
            img_path = include_img
            if os.path.exists(img_path):
                pic = slide.shapes.add_picture(img_path, x, current_y, width=width)
                current_y += pic.height + Cm(1)
            else:
                placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, current_y, width, Cm(12))
                placeholder.text_frame.text = f"[Insert {img_path}]"
                current_y += Cm(13)
                
            if img_caption:
                cap_box = slide.shapes.add_textbox(x, current_y, width, Cm(3))
                tf_cap = cap_box.text_frame
                tf_cap.word_wrap = True
                p_cap = tf_cap.paragraphs[0]
                p_cap.text = img_caption
                p_cap.font.italic = True
                p_cap.font.size = Pt(22)
                p_cap.font.color.rgb = RGBColor(85, 85, 85)
                
                num_cap_chars = len(img_caption)
                cap_lines = math.ceil(num_cap_chars / 45.0)
                current_y += Cm(cap_lines * 1.2 + 2)
                
        return current_y + Cm(2) # padding before next section

    # Columns Setup
    col1_x = Cm(4)
    col2_x = Cm(30.5)
    col3_x = Cm(57)
    
    # COLUMN 1
    y1 = Cm(19)
    y1 = add_section(col1_x, y1, "Introduction", "PlannEd is a next-generation Android application designed to eliminate human bias in scheduling. It introduces a totally native, offline Multi-Layer Perceptron (MLP) neural network using Jetpack Compose to calculate user-specific task padding dynamically. This is coupled with offline Natural Language Processing (NLP) for intuitive hands-free Voice control, allowing fluid schedule management without network latency.")
    
    y1 = add_section(col1_x, y1, "Abstract", "Traditional planners require manual configuration of buffer times and tedious typing. Users often under-estimate the time tasks take, leading to missed deadlines and burnout. PlannEd solves this by synthesizing an offline deep neural network directly into Kotlin matrices, applying dynamic baseline padding to tasks while supporting seamless hands-free operation via a custom native NLP engine.")
    
    y1 = add_section(col1_x, y1, "Problem Statement", "Scheduling apps struggle with two friction points: (1) Manual data entry is slow. (2) Humans naturally underestimate task durations, resulting in packed schedules. Creating an intelligent padding assistant usually requires a pinging out to cloud-based LLMs or hosting heavy offline TFLite libraries, both of which severely bloat the application memory and introduce latency.")

    # COLUMN 2
    y2 = Cm(19)
    y2 = add_section(col2_x, y2, "Literature Review", "Current productivity apps utilize either fixed manual buffers or heavyweight API integrations (e.g. OpenAI). For offline modeling, TensorFlow Lite is an industry standard but introduces ~30MB of library overhead. Studies on the Planning Fallacy suggest that dynamic padding predictions based on regression models of past human behavior significantly optimize realistic completion probability.")
    
    y2 = add_section(col2_x, y2, "Methodology", "Architecture Components\nOffline Model Training: A neural network was trained externally using PyTorch over 300 epochs to map task durations to human padding requirements.\n\nThe Bridge: Instead of adding heavy dependencies, a custom bridging layer extracted weights (W1, b1, W2, b2) into JSON. Inside the Android runtime, pure Kotlin matrix multiplication executes the forward pass natively at zero-latency.")
    
    y2 = add_section(col2_x, y2, "Training & Validation", "", include_img="model_loss.png", img_caption="The model converged safely over 300 epochs. The residuals displayed a tight normal distribution clustered around zero, resulting in a Mean Absolute Error (MAE) of merely 1.77 minutes, proving gradient descent efficiently navigated local minimums.")

    # COLUMN 3
    y3 = Cm(19)
    y3 = add_section(col3_x, y3, "Results & Graphs", "The scatter predictions show a tight correlation between the initial duration and expected padding limits. Residual distributions form a steep normal curve centered securely at zero, proving the optimizer navigated the local minimums perfectly without detrimental overfitting.", include_img="graph_predictions.png")
    
    y3 = add_section(col3_x, y3, "More Results", "", include_img="graph_residuals.png")
    
    y3 = add_section(col3_x, y3, "Conclusion", "The custom ML framework enables a blazingly fast inference engine that solves the scheduling padding projection problem offline.")
    
    y3 = add_section(col3_x, y3, "Future Work", "Future iterations will fully unleash the Adaptive Time Intelligence (ATI) bridge. As a user completes tasks natively, the framework will blend the pre-trained neural network baseline with the user's localized regression model (in 10% steps), infinitely adapting the padding to the individual's unique procrastination behaviors over time.")
    
    y3 = add_section(col3_x, y3, "References", "1. PlannEd App Analytics and Design Documentation.\n2. The Planning Fallacy - Psychological studies on time estimation.\n3. PyTorch & NumPy Documentation - Lightweight native exports.\n4. Scikit-learn documentation for MLP concepts.")

    prs.save('PlannEd_A0_Premium_Poster_v2.pptx')
    print("PPTX generated successfully as PlannEd_A0_Premium_Poster_v2.pptx")

if __name__ == "__main__":
    create_poster()
